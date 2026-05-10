#!/usr/bin/env python3
"""Generate capstone presentation and printable packet files for a build packet.

Outputs:
  - capstone-deck.md      (always)
  - capstone-deck.pptx    (when python-pptx is available)
  - print-packet.md       (always)
  - print-packet.html     (always)
  - print-packet.pdf      (when reportlab or weasyprint is available)
  - capstone-manifest.json

v4.2: asset discovery supports root-mode packets (`drawings/`, `images/`),
post-reorg packets (`build/packet` with sibling `build/drawings`), repo-level
`assets/images`, and explicit `--drawings-dir` / `--images-dir` overrides.
"""

from __future__ import annotations

import argparse
import csv
import datetime as dt
import html
import json
import os
import re
from pathlib import Path


FILE_PURPOSES = {
    "design.md": "Project intent, catalog metadata, assumptions, and validation plan.",
    "bom.csv": "Starter bill of materials with part categories, quantities, drawing refs, and notes.",
    "sourcing.csv": "Supplier/search tracker with specs, price/date fields, lead time, substitutes, and risks.",
    "cut-list.csv": "Rough/final stock sizes, material, grain/orientation, operations, yield, and offcuts.",
    "validation.csv": "Target/measured values, tolerance, environment, result, and tuning/build action log.",
    "assembly-manual.md": "Shop-facing sequence, tools, fixtures, safety, tuning, finishing, and maintenance notes.",
    "supplier-rfq.md": "Supplier email/request-for-quote starter.",
    "visual-bom-brief.md": "Art direction for an image-forward visual BOM.",
    "drawing-brief.md": "Manufacturing drawing and technical product sketch brief.",
    "wolfram-starter.wl": "Wolfram starter for physics, optimization, visualization, and validation.",
}

PACKET_ORDER = [
    "design.md",
    "bom.csv",
    "sourcing.csv",
    "cut-list.csv",
    "drawing-brief.md",
    "assembly-manual.md",
    "validation.csv",
    "supplier-rfq.md",
    "visual-bom-brief.md",
    "wolfram-starter.wl",
]

# Files this script *generates* — must never be treated as packet inputs on
# subsequent regenerations.  Without this filter, regenerating folds the
# previous capstone-deck.md and print-packet.md back into the next packet
# run, causing recursive content growth in the print packet and stale
# duplicate sections in the deck.  Surfaced during the moseno packet regen;
# fix shipped in v3.1.
GENERATED_OUTPUTS = {
    "capstone-deck.md",
    "capstone-deck.pptx",
    "print-packet.md",
    "print-packet.html",
    "print-packet.pdf",
    "capstone-manifest.json",
}


def packet_files(packet_dir: Path) -> list[Path]:
    files = []
    for name in PACKET_ORDER:
        path = packet_dir / name
        if path.exists():
            files.append(path)
    for path in sorted(packet_dir.iterdir()):
        if (
            path.is_file()
            and path.name not in PACKET_ORDER
            and path.name not in GENERATED_OUTPUTS  # v3.1: never re-fold our own outputs
            and not path.name.startswith(".")        # v3.1: skip dotfiles (.v3-issue-body.md, etc.)
            and path.suffix.lower() in {".md", ".csv", ".wl", ".txt"}
        ):
            files.append(path)
    return files


def title_from_packet(packet_dir: Path, fallback: str | None) -> str:
    if fallback:
        return fallback
    design = packet_dir / "design.md"
    if design.exists():
        for line in design.read_text(encoding="utf-8", errors="replace").splitlines():
            if line.startswith("# "):
                return line[2:].strip()
    return packet_dir.name.replace("-", " ").title()


def read_csv_as_markdown(path: Path, max_rows: int | None = None) -> str:
    rows = []
    with path.open(newline="", encoding="utf-8", errors="replace") as handle:
        for row in csv.reader(handle):
            rows.append(row)
            if max_rows is not None and len(rows) > max_rows:
                break
    if not rows:
        return ""
    width = max(len(row) for row in rows)
    padded = [row + [""] * (width - len(row)) for row in rows]
    lines = [
        "| " + " | ".join(cell.replace("|", "\\|") for cell in padded[0]) + " |",
        "| " + " | ".join(["---"] * width) + " |",
    ]
    for row in padded[1:]:
        lines.append("| " + " | ".join(cell.replace("|", "\\|") for cell in row) + " |")
    return "\n".join(lines)


def read_text_for_packet(path: Path) -> str:
    if path.suffix.lower() == ".csv":
        return read_csv_as_markdown(path)
    if path.suffix.lower() == ".wl":
        return "```wolfram\n" + path.read_text(encoding="utf-8", errors="replace").strip() + "\n```"
    return path.read_text(encoding="utf-8", errors="replace").strip()


def csv_rows(path: Path) -> list[list[str]]:
    rows: list[list[str]] = []
    with path.open(newline="", encoding="utf-8", errors="replace") as handle:
        for row in csv.reader(handle):
            rows.append(row)
    return rows


def first_lines(path: Path, count: int = 5) -> list[str]:
    if path.suffix.lower() == ".csv":
        text = read_csv_as_markdown(path, max_rows=5)
    else:
        text = path.read_text(encoding="utf-8", errors="replace")
    lines = []
    for line in text.splitlines():
        stripped = re.sub(r"^[#*\-\s|`]+", "", line).strip()
        if stripped and not set(stripped) <= {"-", "|", ":"}:
            lines.append(stripped)
        if len(lines) >= count:
            break
    return lines


IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp", ".svg"}
DRAWING_SUFFIXES = {".svg", ".pdf", ".png", ".jpg", ".jpeg", ".dxf"}


def _unique_existing_dirs(paths: list[Path]) -> list[Path]:
    """Return existing directories once, preserving discovery order."""
    seen: set[Path] = set()
    dirs: list[Path] = []
    for path in paths:
        resolved = path.resolve()
        if resolved in seen or not resolved.exists() or not resolved.is_dir():
            continue
        seen.add(resolved)
        dirs.append(resolved)
    return dirs


def _asset_dirs(packet_dir: Path, kind: str, explicit_dirs: list[Path] | None = None) -> list[Path]:
    """Find asset folders for root-mode and post-reorg packet layouts.

    Root-mode packets keep `drawings/` and `images/` beside `design.md`.
    Reorged repos often call this script from `build/packet/`, while assets
    live in sister folders such as `build/drawings/` or repo-level
    `assets/images/`. Discover common layouts automatically and allow
    explicit folders for unusual repos.
    """
    packet_dir = packet_dir.resolve()
    candidates: list[Path] = []
    if explicit_dirs:
        candidates.extend(explicit_dirs)
    candidates.extend([
        packet_dir / kind,
        packet_dir.parent / kind,
        packet_dir.parent.parent / kind,
    ])
    if kind == "images":
        candidates.extend([
            packet_dir / "assets" / "images",
            packet_dir.parent / "assets" / "images",
            packet_dir.parent.parent / "assets" / "images",
        ])
    return _unique_existing_dirs(candidates)


def _markdown_ref(packet_dir: Path, path: Path) -> str:
    """Return a portable markdown reference from packet_dir to path."""
    return os.path.relpath(path.resolve(), packet_dir.resolve()).replace(os.sep, "/")


def image_files(packet_dir: Path, image_dirs: list[Path] | None = None) -> list[Path]:
    """Return image files from images/, including SVGs.

    v3.1: SVG files in images/ used to be silently dropped, leaving the
    Images slide empty when the only renders shipped were SVG.  We now
    include them; the PPTX writer rasterizes them through `rasterize_svg`
    just like drawings/*.svg.
    """
    images = []
    for image_dir in _asset_dirs(packet_dir, "images", image_dirs):
        for path in sorted(image_dir.rglob("*")):
            # Skip files inside the .thumbnails cache directory we maintain.
            if ".thumbnails" in path.parts:
                continue
            if path.suffix.lower() in IMAGE_SUFFIXES:
                images.append(path)
    return images


def file_map_table(files: list[Path]) -> str:
    rows = ["| File | Purpose |", "| --- | --- |"]
    for path in files:
        purpose = FILE_PURPOSES.get(path.name, "Project artifact.")
        rows.append(f"| `{path.name}` | {purpose} |")
    return "\n".join(rows)


def design_intent_from_packet(packet_dir: Path) -> str:
    """Extract the project/design intent from design.md.

    Tries common heading variants first (`## Intent`, `## Design Intent`,
    `## Project Intent`), then falls back to the first non-empty paragraph
    after the H1.  This matters because the deck's slide 2 is the only place
    instrument-specific intent surfaces; a generic placeholder there is the
    most-noticed quality regression in v1 of this skill.
    """
    design = packet_dir / "design.md"
    if not design.exists():
        return "See design.md for intent and assumptions."
    text = design.read_text(encoding="utf-8", errors="replace")

    # Match `## Intent`, `## Design Intent`, `## Project Intent`, etc.
    match = re.search(
        r"^##\s*(?:Design\s+|Project\s+|Build\s+)?Intent\b\s*\n+([^\n][^\n]*(?:\n(?!##)[^\n]*)*)",
        text,
        re.M,
    )
    if match:
        intent = match.group(1).strip()
        return intent.split("\n\n", 1)[0]

    # Fallback: first non-empty paragraph after the H1.
    after_h1 = re.split(r"^#\s+.+$", text, maxsplit=1, flags=re.M)
    if len(after_h1) == 2:
        for paragraph in re.split(r"\n\s*\n", after_h1[1]):
            stripped = paragraph.strip()
            if stripped and not stripped.startswith(("##", "```", "|", "-")):
                return stripped.split("\n\n", 1)[0]

    return "See design.md for intent."


def governing_model_from_packet(packet_dir: Path) -> tuple[str, list[str]]:
    """Extract (prose_summary, [formula_block, ...]) from design.md.

    Looks for `## Governing Model` (and common variants).  Returns a list of
    formulas — instruments with multiple coupled physics models (the udu's
    dual Helmholtz, the ceramic-electric-violin's string + body coupling)
    deserve to surface ALL fenced blocks on the Physics Model slide, not
    just the first.

    v3.1 fix: previously returned only the first fenced block, which meant
    the Helmholtz equation got dropped on the violin packet after the
    string equation.
    """
    design = packet_dir / "design.md"
    if not design.exists():
        return ("", [])
    text = design.read_text(encoding="utf-8", errors="replace")

    match = re.search(
        r"^##\s*(?:Governing\s+Model|Physics\s+Model|Acoustic\s+Model|Governing\s+Equations?)\b\s*\n+(.*?)(?=^##\s|\Z)",
        text,
        re.M | re.S,
    )
    if not match:
        return ("", [])
    section = match.group(1).strip()

    # Extract ALL fenced code blocks as formulas.  Order is preserved so the
    # primary physics renders first on the slide.
    formulas = [m.strip() for m in re.findall(r"```(?:[a-zA-Z]*)\n(.*?)```", section, re.S)]

    # The prose is whatever precedes the first fenced block.  If there are no
    # fenced blocks, treat the whole section as prose.
    first_fence = re.search(r"```", section)
    prose = section[: first_fence.start()].strip() if first_fence else section.strip()

    # Keep prose to first paragraph only — slides have limited room.
    prose = prose.split("\n\n", 1)[0]
    return (prose, formulas)


def hardware_alignment_from_packet(packet_dir: Path) -> tuple[str, list[list[str]]]:
    """Pull (prose, table_rows) for the Hardware Alignment slide.

    Looks for `## Hardware Alignment`, `## Pipeline`, or `## Hardware Pipeline`
    in design.md.  Returns the intro paragraph AND any markdown table that
    follows it — hybrid instruments (ceramic-electric-violin) describe their
    pipeline as a multi-row table (slip-cast body + wood neck + electronics)
    and dropping the table loses the whole point of the slide.

    v3.1 fix: previously returned only the first paragraph.
    """
    design = packet_dir / "design.md"
    if not design.exists():
        return ("", [])
    text = design.read_text(encoding="utf-8", errors="replace")
    match = re.search(
        r"^##\s*(?:Hardware\s+Alignment|Hardware\s+Pipeline|Pipeline|Shop\s+Pipeline)\b\s*\n+(.*?)(?=^##\s|\Z)",
        text,
        re.M | re.S,
    )
    if not match:
        return ("", [])
    section = match.group(1).strip()

    # Pull a markdown table out of the section if one exists.
    table_rows: list[list[str]] = []
    for line in section.splitlines():
        s = line.strip()
        if s.startswith("|") and "|" in s[1:]:
            # Skip the alignment row (dashes/colons only).
            if set(s.replace("|", "").replace(":", "").replace("-", "").replace(" ", "")) == set():
                continue
            cells = [c.strip() for c in s.strip("|").split("|")]
            table_rows.append(cells)
        elif table_rows:
            # Stop at the first non-table line after we've started collecting.
            break

    # Prose: everything before the first table row, first paragraph only.
    if table_rows:
        first_table_idx = section.find(next(line for line in section.splitlines() if line.strip().startswith("|")))
        prose = section[:first_table_idx].strip()
    else:
        prose = section
    prose = prose.split("\n\n", 1)[0]
    return (prose, table_rows)


def family_spec_from_packet(packet_dir: Path) -> list[list[str]]:
    """Return rows of a family spec table if one exists.

    Looks for either a `family-spec.csv` file in the packet OR a markdown
    table under `## Family Targets` / `## Current Family Targets` in
    design.md.  Returns CSV-style rows (header first); empty list if none
    found.  Drives the Family Spec slide AND the family-aware OpenSCAD
    starter generator.
    """
    csv_path = packet_dir / "family-spec.csv"
    if csv_path.exists():
        rows = csv_rows(csv_path)
        if rows:
            return rows

    design = packet_dir / "design.md"
    if not design.exists():
        return []
    text = design.read_text(encoding="utf-8", errors="replace")
    match = re.search(
        r"^##\s*(?:Current\s+)?Family\s+(?:Targets?|Spec)\b\s*\n+(.*?)(?=^##\s|\Z)",
        text,
        re.M | re.S,
    )
    if not match:
        return []
    section = match.group(1)

    # Extract the first markdown table.
    rows: list[list[str]] = []
    for line in section.splitlines():
        stripped = line.strip()
        if not stripped.startswith("|"):
            if rows:  # table ended
                break
            continue
        # Skip alignment row like |---|---|
        if set(stripped.replace("|", "").replace(":", "").replace("-", "").replace(" ", "")) == set():
            continue
        cells = [c.strip() for c in stripped.strip("|").split("|")]
        rows.append(cells)
    return rows


def drawing_files(packet_dir: Path, drawing_dirs: list[Path] | None = None) -> list[Path]:
    """Return drawings from root-mode or sibling drawings folders."""
    drawings = []
    for drawing_dir in _asset_dirs(packet_dir, "drawings", drawing_dirs):
        for path in sorted(drawing_dir.rglob("*")):
            if path.suffix.lower() in DRAWING_SUFFIXES:
                drawings.append(path)
    return drawings


def rasterize_svg(svg_path: Path, png_path: Path, width: int = 800) -> bool:
    """Best-effort SVG -> PNG conversion for thumbnail rendering.

    Tries cairosvg, then ImageMagick `convert`, then `rsvg-convert`.  Returns
    True if any of them succeeded.  Used by the Drawing References slide so
    drawings/*.svg can be embedded as thumbnails alongside the bitmap images
    in images/.
    """
    if png_path.exists() and png_path.stat().st_mtime >= svg_path.stat().st_mtime:
        return True
    try:
        import cairosvg  # type: ignore
        cairosvg.svg2png(url=str(svg_path), write_to=str(png_path), output_width=width)
        return True
    except Exception:
        pass
    import shutil
    import subprocess
    for cmd in (
        ["convert", "-density", "120", str(svg_path), str(png_path)],
        ["rsvg-convert", "-w", str(width), str(svg_path), "-o", str(png_path)],
    ):
        if shutil.which(cmd[0]):
            try:
                subprocess.run(cmd, check=True, capture_output=True)
                return True
            except Exception:
                continue
    return False


def build_deck_slides(
    packet_dir: Path,
    title: str,
    files: list[Path],
    image_dirs: list[Path] | None = None,
    drawing_dirs: list[Path] | None = None,
) -> list[dict[str, object]]:
    """Build the slide list, with content-aware slides where data is available.

    Slide order:
      1. Title
      2. Project Intent (from design.md `## Intent` / `## Design Intent`)
      3. Physics Model (when `## Governing Model` exists in design.md)
      4. Hardware Alignment (when `## Hardware Alignment` exists in design.md)
      5. How To Use This Packet
      6. File Map
      7. Family Spec (when family-spec.csv or family table is present)
      8. Build Workflow
      9. Sourcing And BOM
     10. Shop Packet
     11. Drawing References (when drawings/ has SVGs/PDFs)
     12. Images And Screenshots (with embedded thumbnails when present)
     13. Validation Plan
     14. Open Risks / Decisions
     15. Next Actions

    Slides 3, 4, 7, 11 are conditional — they only appear when the
    underlying data exists.  This keeps the deck honest: an instrument with
    no family doesn't get an empty Family Spec slide saying "TBD"; a project
    with no governing-model section doesn't get a placeholder physics slide.
    """
    images = image_files(packet_dir, image_dirs)
    drawings = drawing_files(packet_dir, drawing_dirs)
    intent = design_intent_from_packet(packet_dir)
    physics_prose, physics_formulas = governing_model_from_packet(packet_dir)
    hardware_prose, hardware_table = hardware_alignment_from_packet(packet_dir)
    family_rows = family_spec_from_packet(packet_dir)

    slides: list[dict[str, object]] = []

    # 1. Title
    slides.append({
        "title": title,
        "bullets": [
            "Musical instrument documentation capstone",
            f"Build packet: {packet_dir.name}",
            f"Generated: {dt.date.today().isoformat()}",
        ],
    })

    # 2. Project Intent (always present; falls back gracefully)
    slides.append({
        "title": "Project Intent",
        "bullets": [intent],
        "speaker_notes": "Read design.md before committing to dimensions or sourcing decisions.",
    })

    # 3. Physics Model (conditional) — renders ALL fenced blocks from
    # `## Governing Model`, not just the first.  Coupled-physics instruments
    # (udu dual-Helmholtz, ceramic-electric-violin string+body) need every
    # equation visible on the slide.
    if physics_prose or physics_formulas:
        physics_bullets: list = []
        if physics_prose:
            physics_bullets.append(physics_prose)
        for formula in physics_formulas:
            physics_bullets.append({"kind": "code", "text": formula})
        slides.append({
            "title": "Physics Model",
            "bullets": physics_bullets,
            "speaker_notes": (
                "Governing equations extracted verbatim from design.md. "
                "Apply empirical corrections (NAF K2, scale offsets) only where the model permits — see references/acoustic-models.md."
            ),
        })

    # 4. Hardware Alignment (conditional) — renders both the intro paragraph
    # AND any markdown table under the heading.  Hybrid instruments use the
    # table to describe multiple sub-pipelines (slip-cast body + wood neck +
    # electronics on the ceramic-electric-violin); dropping the table loses
    # the slide's purpose.
    if hardware_prose or hardware_table:
        slide: dict[str, object] = {
            "title": "Hardware Alignment",
            "bullets": [hardware_prose] if hardware_prose else [],
            "speaker_notes": (
                "Identifies which shop pipeline(s) this instrument lives in: Bambu+kiln slip-cast, "
                "40W laser flat-pack, CNC+lathe, segmented turning, drum-skin work, or hybrid combinations."
            ),
        }
        if hardware_table and len(hardware_table) >= 2:
            slide["table"] = hardware_table
        slides.append(slide)

    # 5. How To Use This Packet
    slides.append({
        "title": "How To Use This Packet",
        "bullets": [
            "Start with design.md for intent and assumptions.",
            "Use bom.csv, sourcing.csv, and cut-list.csv before buying or cutting.",
            "Use drawing-brief.md and CAD/CNC folders before machining.",
            "Print the packet for shopping, shop work, and validation.",
        ],
    })

    # 6. File Map
    slides.append({
        "title": "File Map",
        "bullets": [f"{path.name}: {FILE_PURPOSES.get(path.name, 'Project artifact.')}" for path in files[:8]],
    })

    # 7. Family Spec (conditional)
    if family_rows and len(family_rows) > 1:  # need header + at least one variant
        slides.append({
            "title": "Family Spec",
            "table": family_rows,
            "speaker_notes": (
                "Sizes scale via the master scale factor; tuning targets are first-order Helmholtz/cantilever "
                "predictions to be empirically corrected per prototype."
            ),
        })

    # 8. Build Workflow
    slides.append({
        "title": "Build Workflow",
        "bullets": [
            "Design and assumptions",
            "Source materials and hardware",
            "Prepare stock, fixtures, and CNC/laser/lathe setup",
            "Assemble, tune, finish, and validate",
        ],
    })

    # 9. Sourcing And BOM
    slides.append({
        "title": "Sourcing And BOM",
        "bullets": [
            "BOM gives part categories and drawing references.",
            "Sourcing tracks search terms, supplier candidates, price/date, lead time, substitutions.",
            "Visual BOM brief turns the parts list into a presentation-ready image board.",
        ],
    })

    # 10. Shop Packet
    slides.append({
        "title": "Shop Packet",
        "bullets": [
            "Cut list for lumber/sheet/blank planning.",
            "Assembly manual for away-from-keyboard work.",
            "Validation sheet for measured dimensions, tuning, pass/fail checks.",
        ],
    })

    # 11. Drawing References (conditional, with thumbnails)
    if drawings:
        slides.append({
            "title": "Drawings, CAD, CNC",
            "bullets": [
                "drawing-brief.md defines required views, dimensions, datums, sketch intent.",
                "cad/ holds models and design tables.",
                "cnc/ holds CAM, toolpaths, setup sheets, dry-run notes.",
                "drawings/ holds PDFs, SVGs, DXFs, drawing exports.",
            ],
            "drawing_thumbnails": [_markdown_ref(packet_dir, d) for d in drawings[:4]],
        })
    else:
        slides.append({
            "title": "Drawings, CAD, CNC",
            "bullets": [
                "drawing-brief.md defines required views, dimensions, datums, sketch intent.",
                "cad/ holds models and design tables.",
                "cnc/ holds CAM, toolpaths, setup sheets, dry-run notes.",
                "drawings/ holds PDFs, SVGs, DXFs, drawing exports.",
            ],
        })

    # 12. Images And Screenshots — embed thumbnails when present
    if images:
        slides.append({
            "title": "Images And Screenshots",
            "bullets": [_markdown_ref(packet_dir, path) for path in images[:6]],
            "image_thumbnails": [_markdown_ref(packet_dir, p) for p in images[:6]],
        })
    else:
        slides.append({
            "title": "Images And Screenshots",
            "bullets": [
                "Add hero render/photo, visual BOM, shop screenshots, drawing previews, validation photos in images/.",
            ],
        })

    # 13. Validation Plan
    slides.append({
        "title": "Validation Plan",
        "bullets": [
            "A4 = 440 Hz reference check.",
            "Tuning targets logged in validation.csv.",
            "Critical dimensions verified against design sheet and CAD.",
            "Photos and revision notes after each major step.",
        ],
    })

    # 14. Open Risks / Decisions
    slides.append({
        "title": "Open Risks / Decisions",
        "bullets": [
            "TBDs in design sheet and BOM.",
            "Supplier price/availability not yet verified.",
            "Generated images marked as concept placeholders.",
            "Empirical corrections await measured prototype data.",
        ],
    })

    # 15. Next Actions
    slides.append({
        "title": "Next Actions",
        "bullets": [
            "Replace TBDs with measured/source-backed values.",
            "Verify live supplier price and availability before buying.",
            "Export final drawings and visual BOM images.",
            "Regenerate this deck and print packet after final edits.",
        ],
    })

    return slides


def write_deck_markdown(path: Path, slides: list[dict[str, object]]) -> None:
    """Write the markdown source for the deck, mirroring slide structure 1:1.

    Bullets can be strings (plain) or dicts (e.g., `{"kind": "code", "text": ...}`
    for the Physics Model formula).  Tables (Family Spec) and thumbnail
    references (drawings, images) render as markdown tables and image links.
    """
    parts = []
    for slide in slides:
        parts.append(f"# {slide['title']}")
        for bullet in slide.get("bullets", []):
            if isinstance(bullet, dict) and bullet.get("kind") == "code":
                parts.append("")
                parts.append("```")
                parts.append(str(bullet.get("text", "")))
                parts.append("```")
            else:
                parts.append(f"- {bullet}")
        # Family Spec / other slides with a structured table
        table = slide.get("table")
        if table:
            parts.append("")
            width = max(len(row) for row in table)
            padded = [row + [""] * (width - len(row)) for row in table]
            parts.append("| " + " | ".join(c.replace("|", "\\|") for c in padded[0]) + " |")
            parts.append("| " + " | ".join(["---"] * width) + " |")
            for row in padded[1:]:
                parts.append("| " + " | ".join(c.replace("|", "\\|") for c in row) + " |")
        # Inline thumbnail references for drawings/images slides
        for ref_key, label in (("drawing_thumbnails", "drawing"), ("image_thumbnails", "image")):
            refs = slide.get(ref_key) or []
            if refs:
                parts.append("")
                for ref in refs:
                    parts.append(f"![{label}]({ref})")
        notes = slide.get("speaker_notes")
        if notes:
            parts.append("")
            parts.append(f"_Speaker notes:_ {notes}")
        parts.append("")
        parts.append("---")
        parts.append("")
    path.write_text("\n".join(parts).strip() + "\n", encoding="utf-8")


def write_print_packet(path: Path, packet_dir: Path, title: str, files: list[Path]) -> None:
    parts = [
        f"# {title} Print Packet",
        "",
        f"Generated: {dt.date.today().isoformat()}",
        f"Packet folder: `{packet_dir}`",
        "",
        "## File Map",
        "",
        file_map_table(files),
        "",
    ]
    for file_path in files:
        purpose = FILE_PURPOSES.get(file_path.name, "Project artifact.")
        parts.extend(
            [
                '<div class="page-break"></div>',
                "",
                f"## {file_path.name}",
                "",
                purpose,
                "",
                read_text_for_packet(file_path),
                "",
            ]
        )
    path.write_text("\n".join(parts).strip() + "\n", encoding="utf-8")


def markdown_to_html(markdown: str, title: str) -> str:
    lines = markdown.splitlines()
    out = []
    in_code = False
    in_ul = False
    in_table = False

    def close_lists() -> None:
        nonlocal in_ul
        if in_ul:
            out.append("</ul>")
            in_ul = False

    for line in lines:
        if line.startswith("```"):
            close_lists()
            if not in_code:
                out.append("<pre><code>")
                in_code = True
            else:
                out.append("</code></pre>")
                in_code = False
            continue
        if in_code:
            out.append(html.escape(line))
            continue
        if line.startswith('<div class="page-break"></div>'):
            close_lists()
            out.append('<div class="page-break"></div>')
            continue
        if line.startswith("| ") and line.endswith(" |"):
            close_lists()
            cells = [html.escape(cell.strip()) for cell in line.strip("|").split("|")]
            if set(cell.strip("-: ") for cell in cells) == {""}:
                continue
            if not in_table:
                out.append("<table>")
                in_table = True
            tag = "th" if not any("<tr>" in row for row in out[-2:]) else "td"
            out.append("<tr>" + "".join(f"<{tag}>{cell}</{tag}>" for cell in cells) + "</tr>")
            continue
        if in_table:
            out.append("</table>")
            in_table = False
        if not line.strip():
            close_lists()
            out.append("")
        elif line.startswith("# "):
            close_lists()
            out.append(f"<h1>{html.escape(line[2:].strip())}</h1>")
        elif line.startswith("## "):
            close_lists()
            out.append(f"<h2>{html.escape(line[3:].strip())}</h2>")
        elif line.startswith("### "):
            close_lists()
            out.append(f"<h3>{html.escape(line[4:].strip())}</h3>")
        elif line.startswith("- "):
            if not in_ul:
                out.append("<ul>")
                in_ul = True
            out.append(f"<li>{html.escape(line[2:].strip())}</li>")
        else:
            close_lists()
            out.append(f"<p>{html.escape(line)}</p>")
    close_lists()
    if in_table:
        out.append("</table>")

    css = """
body { font-family: Arial, sans-serif; color: #111; line-height: 1.35; max-width: 900px; margin: 32px auto; }
h1 { font-size: 28px; border-bottom: 2px solid #222; padding-bottom: 6px; }
h2 { font-size: 22px; margin-top: 28px; }
h3 { font-size: 18px; }
table { border-collapse: collapse; width: 100%; margin: 12px 0 20px; font-size: 12px; }
th, td { border: 1px solid #888; padding: 6px; vertical-align: top; }
th { background: #e8eef8; }
pre { white-space: pre-wrap; background: #f5f5f5; padding: 10px; border: 1px solid #ddd; }
.page-break { page-break-before: always; break-before: page; }
@media print { body { margin: 0.5in; max-width: none; } a { color: #111; } }
"""
    return f"<!doctype html><html><head><meta charset=\"utf-8\"><title>{html.escape(title)}</title><style>{css}</style></head><body>\n" + "\n".join(out) + "\n</body></html>\n"


def write_pptx_if_available(path: Path, slides: list[dict[str, object]], title: str, packet_dir: Path | None = None) -> bool:
    """Render the deck as a styled .pptx.

    Handles four slide content types:
      - plain bullets (str)
      - code-style bullets (dict with kind="code") — rendered monospace
      - tables (slide["table"] = list[list[str]]) — rendered as a styled grid
      - thumbnail rows (slide["image_thumbnails"] / "drawing_thumbnails") —
        bitmap PNGs / SVGs are rasterized and embedded with `add_picture`,
        captioned below

    `packet_dir` is needed to resolve relative thumbnail paths from the
    slide dicts.
    """
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches, Pt, Emu  # type: ignore
        from pptx.dml.color import RGBColor  # type: ignore
        from pptx.enum.shapes import MSO_SHAPE  # type: ignore
    except Exception:
        return False

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    accent = RGBColor(0x1F, 0x3A, 0x68)
    accent_light = RGBColor(0xE8, 0xEE, 0xF8)
    text_color = RGBColor(0x11, 0x11, 0x11)
    sub_color = RGBColor(0x44, 0x55, 0x77)

    for index, slide_data in enumerate(slides):
        layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(layout)

        # Top accent band
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.45))
        band.line.fill.background()
        band.fill.solid()
        band.fill.fore_color.rgb = accent

        # Title text
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(12.3), Inches(1.0))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_para = title_tf.paragraphs[0]
        title_para.text = str(slide_data["title"])
        title_para.font.size = Pt(36) if index == 0 else Pt(30)
        title_para.font.bold = True
        title_para.font.color.rgb = accent

        # Subtitle / bullets — split between text body region and thumbnail region.
        bullets = slide_data.get("bullets", [])
        thumbnails = slide_data.get("image_thumbnails") or slide_data.get("drawing_thumbnails") or []
        has_thumbs = bool(thumbnails) and packet_dir is not None
        has_table = bool(slide_data.get("table"))

        # Text body shrinks horizontally when thumbnails are present so they fit beside it.
        body_top = Inches(1.7) if index == 0 else Inches(1.6)
        body_width = Inches(6.5) if has_thumbs else Inches(12.3)
        body_box = slide.shapes.add_textbox(Inches(0.5), body_top, body_width, Inches(5.2))
        body_tf = body_box.text_frame
        body_tf.word_wrap = True

        first_para_used = False
        for bullet in bullets:
            para = body_tf.paragraphs[0] if not first_para_used else body_tf.add_paragraph()
            first_para_used = True
            if isinstance(bullet, dict) and bullet.get("kind") == "code":
                # Render formula in a monospace, slightly smaller font.
                para.text = str(bullet.get("text", ""))
                para.font.name = "Consolas"
                para.font.size = Pt(16)
                para.font.color.rgb = accent
                para.space_after = Pt(10)
            else:
                para.text = ("• " if index != 0 else "") + str(bullet)
                para.font.size = Pt(22) if index == 0 else Pt(20)
                para.font.color.rgb = text_color if index == 0 else sub_color
                para.space_after = Pt(8)

        # Render a table for slides with structured tabular data (Family Spec).
        if has_table:
            table_rows = slide_data["table"]
            n_rows = len(table_rows)
            n_cols = max(len(row) for row in table_rows)
            table_left = Inches(0.5)
            table_top = body_top + Inches(0.2 + 0.3 * len(bullets))
            table_width = Inches(12.3)
            table_height = Inches(min(4.5, 0.4 + 0.35 * n_rows))
            tbl_shape = slide.shapes.add_table(n_rows, n_cols, table_left, table_top, table_width, table_height)
            tbl = tbl_shape.table
            for ri, row in enumerate(table_rows):
                padded = list(row) + [""] * (n_cols - len(row))
                for ci, cell_text in enumerate(padded):
                    cell = tbl.cell(ri, ci)
                    cell.text = cell_text
                    para = cell.text_frame.paragraphs[0]
                    para.font.size = Pt(13) if ri == 0 else Pt(12)
                    para.font.bold = (ri == 0)
                    para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if ri == 0 else text_color
                    cell.fill.solid()
                    if ri == 0:
                        cell.fill.fore_color.rgb = accent
                    else:
                        cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if ri % 2 == 1 else accent_light

        # Embed thumbnails on the right side of the slide.
        if has_thumbs:
            thumb_area_left = Inches(7.3)
            thumb_area_top = body_top
            thumb_area_width = Inches(5.5)
            thumb_area_height = Inches(5.2)

            # Resolve and rasterize SVGs; collect bitmap paths.
            bitmap_paths: list[Path] = []
            for ref in thumbnails:
                resolved = (packet_dir / ref).resolve()
                if not resolved.exists():
                    continue
                if resolved.suffix.lower() == ".svg":
                    cache_dir = packet_dir / ".thumbnails"
                    cache_dir.mkdir(exist_ok=True)
                    png_cache = cache_dir / (resolved.stem + ".png")
                    if rasterize_svg(resolved, png_cache, width=900):
                        bitmap_paths.append(png_cache)
                elif resolved.suffix.lower() in {".png", ".jpg", ".jpeg", ".webp"}:
                    bitmap_paths.append(resolved)

            # Layout: 2 columns x 2 rows for up to 4, or stacked for fewer.
            n = len(bitmap_paths)
            if n == 0:
                pass
            elif n == 1:
                cell_w = thumb_area_width
                cell_h = thumb_area_height
                _add_thumbnail(slide, bitmap_paths[0], thumb_area_left, thumb_area_top, cell_w, cell_h, sub_color, Pt)
            else:
                cols = 2 if n > 1 else 1
                rows = (n + cols - 1) // cols
                cell_w = Emu(thumb_area_width // cols)
                cell_h = Emu(thumb_area_height // rows)
                for i, bmp in enumerate(bitmap_paths):
                    r = i // cols
                    c = i % cols
                    left = thumb_area_left + Emu(c * cell_w)
                    top = thumb_area_top + Emu(r * cell_h)
                    _add_thumbnail(slide, bmp, left, top, cell_w, cell_h, sub_color, Pt)

        # Footer
        footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3))
        footer.text_frame.text = f"{title}  |  Slide {index + 1} of {len(slides)}"
        footer.text_frame.paragraphs[0].font.size = Pt(10)
        footer.text_frame.paragraphs[0].font.color.rgb = sub_color

        # Speaker notes
        notes = slide_data.get("speaker_notes")
        if notes and slide.has_notes_slide:
            slide.notes_slide.notes_text_frame.text = str(notes)

    prs.save(path)
    return True


def _add_thumbnail(slide, bitmap_path: Path, left, top, cell_w, cell_h, caption_color, Pt):
    """Add an image with a small caption, scaled to fit cell_w x cell_h with margin."""
    margin_h = int(cell_h * 0.25)  # leave room for caption below
    image_h = cell_h - margin_h
    image_w = cell_w - int(cell_w * 0.05)
    try:
        slide.shapes.add_picture(str(bitmap_path), left + int(cell_w * 0.025), top, width=image_w, height=image_h)
    except Exception:
        return
    # Caption beneath
    cap_box = slide.shapes.add_textbox(left, top + image_h, cell_w, margin_h)
    cap_box.text_frame.word_wrap = True
    cap_para = cap_box.text_frame.paragraphs[0]
    cap_para.text = bitmap_path.name
    cap_para.font.size = Pt(9)
    cap_para.font.color.rgb = caption_color


def write_pdf_if_available(path: Path, packet_dir: Path, title: str, files: list[Path]) -> bool:
    """Render a print-ready PDF directly with reportlab.

    Falls back to weasyprint if reportlab is unavailable. Returns True if a PDF was written.
    """
    try:
        from reportlab.lib.pagesizes import LETTER  # type: ignore
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle  # type: ignore
        from reportlab.lib.units import inch  # type: ignore
        from reportlab.lib import colors  # type: ignore
        from reportlab.platypus import (  # type: ignore
            BaseDocTemplate,
            PageTemplate,
            Frame,
            Paragraph,
            Spacer,
            Table,
            TableStyle,
            PageBreak,
            Preformatted,
        )
        from reportlab.platypus.flowables import KeepTogether  # type: ignore
    except Exception:
        return _write_pdf_via_weasyprint(path, packet_dir, title)

    styles = getSampleStyleSheet()
    h1 = ParagraphStyle("H1", parent=styles["Heading1"], fontSize=18, leading=22, textColor=colors.HexColor("#1F3A68"), spaceAfter=10)
    h2 = ParagraphStyle("H2", parent=styles["Heading2"], fontSize=14, leading=18, textColor=colors.HexColor("#1F3A68"), spaceBefore=12, spaceAfter=6)
    body = ParagraphStyle("Body", parent=styles["BodyText"], fontSize=10, leading=13, spaceAfter=4)
    cover = ParagraphStyle("Cover", parent=styles["Title"], fontSize=24, leading=28, textColor=colors.HexColor("#1F3A68"), alignment=0, spaceAfter=18)
    meta = ParagraphStyle("Meta", parent=styles["BodyText"], fontSize=10, leading=13, textColor=colors.HexColor("#445577"), spaceAfter=4)

    def page_decoration(canvas, doc):  # noqa: ANN001
        canvas.saveState()
        canvas.setFillColor(colors.HexColor("#1F3A68"))
        canvas.rect(0, LETTER[1] - 0.25 * inch, LETTER[0], 0.25 * inch, stroke=0, fill=1)
        canvas.setFillColor(colors.HexColor("#445577"))
        canvas.setFont("Helvetica", 8)
        canvas.drawString(0.5 * inch, 0.4 * inch, f"{title}  |  Generated {dt.date.today().isoformat()}")
        canvas.drawRightString(LETTER[0] - 0.5 * inch, 0.4 * inch, f"Page {doc.page}")
        canvas.restoreState()

    doc = BaseDocTemplate(
        str(path),
        pagesize=LETTER,
        leftMargin=0.6 * inch,
        rightMargin=0.6 * inch,
        topMargin=0.7 * inch,
        bottomMargin=0.7 * inch,
        title=title,
        author="instrument-maker",
    )
    frame = Frame(doc.leftMargin, doc.bottomMargin, doc.width, doc.height, id="normal")
    doc.addPageTemplates([PageTemplate(id="main", frames=frame, onPage=page_decoration)])

    story: list = []

    # --- Cover ---
    story.append(Paragraph(title, cover))
    story.append(Paragraph("Instrument-Maker Print Packet", h2))
    story.append(Spacer(1, 0.1 * inch))
    story.append(Paragraph(f"Build packet folder: <font color='#222'>{html.escape(str(packet_dir))}</font>", meta))
    story.append(Paragraph(f"Generated: {dt.date.today().isoformat()}", meta))
    story.append(Spacer(1, 0.3 * inch))
    story.append(Paragraph("This packet is the printable companion to the build folder. Take it shopping or into the shop. Tear sheets at page breaks.", body))
    story.append(Spacer(1, 0.2 * inch))

    # --- File map ---
    story.append(Paragraph("File Map", h2))
    file_map_data = [["File", "Purpose"]]
    for fp in files:
        file_map_data.append([fp.name, FILE_PURPOSES.get(fp.name, "Project artifact.")])
    fm_table = Table(file_map_data, colWidths=[1.7 * inch, 5.2 * inch])
    fm_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1F3A68")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONT", (0, 0), (-1, 0), "Helvetica-Bold", 10),
        ("FONT", (0, 1), (-1, -1), "Helvetica", 9),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F1F4FA")]),
        ("BOX", (0, 0), (-1, -1), 0.5, colors.HexColor("#444")),
        ("INNERGRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#888")),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 4),
        ("RIGHTPADDING", (0, 0), (-1, -1), 4),
    ]))
    story.append(fm_table)

    # --- One section per file ---
    for fp in files:
        story.append(PageBreak())
        story.append(Paragraph(fp.name, h1))
        purpose = FILE_PURPOSES.get(fp.name, "Project artifact.")
        story.append(Paragraph(f"<i>{html.escape(purpose)}</i>", body))
        story.append(Spacer(1, 0.1 * inch))

        if fp.suffix.lower() == ".csv":
            rows = csv_rows(fp)
            if rows:
                width = max(len(r) for r in rows)
                rows = [r + [""] * (width - len(r)) for r in rows]
                # Truncate cells to keep table reasonable
                rows = [[cell[:120] for cell in row] for row in rows]
                col_count = len(rows[0])
                col_width = (doc.width / max(col_count, 1))
                table = Table(rows, colWidths=[col_width] * col_count, repeatRows=1)
                table.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1F3A68")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                    ("FONT", (0, 0), (-1, 0), "Helvetica-Bold", 8),
                    ("FONT", (0, 1), (-1, -1), "Helvetica", 7),
                    ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F1F4FA")]),
                    ("BOX", (0, 0), (-1, -1), 0.4, colors.HexColor("#444")),
                    ("INNERGRID", (0, 0), (-1, -1), 0.2, colors.HexColor("#AAA")),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("LEFTPADDING", (0, 0), (-1, -1), 2),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 2),
                ]))
                story.append(table)
            else:
                story.append(Paragraph("(empty CSV)", body))
        else:
            text = fp.read_text(encoding="utf-8", errors="replace")
            for raw in text.splitlines():
                line = raw.rstrip()
                if not line.strip():
                    story.append(Spacer(1, 0.06 * inch))
                    continue
                if line.startswith("# "):
                    story.append(Paragraph(html.escape(line[2:].strip()), h1))
                elif line.startswith("## "):
                    story.append(Paragraph(html.escape(line[3:].strip()), h2))
                elif line.startswith("### "):
                    story.append(Paragraph(f"<b>{html.escape(line[4:].strip())}</b>", body))
                elif line.startswith("- ") or line.startswith("* "):
                    story.append(Paragraph(f"• {html.escape(line[2:].strip())}", body))
                elif line.startswith("```"):
                    continue
                else:
                    # Render as paragraph; preserve simple formatting via escape
                    story.append(Paragraph(html.escape(line), body))

    doc.build(story)
    return True


def _write_pdf_via_weasyprint(path: Path, packet_dir: Path, title: str) -> bool:
    try:
        from weasyprint import HTML  # type: ignore
    except Exception:
        return False
    html_path = path.with_suffix(".html")
    if not html_path.exists():
        return False
    HTML(string=html_path.read_text(encoding="utf-8")).write_pdf(str(path))
    return True


VALID_DISPATCH_MODES = {
    "validation-only",
    "direct-script-regen",
    "specialist-dispatched-regen",
    "full-agent-generation",
}


def _resolve_dispatch_mode(arg_value: str | None) -> str | None:
    """Pull dispatch mode from CLI flag or the INSTRUMENT_MAKER_DISPATCH_MODE
    env var. Used to stamp ``capstone-manifest.json`` with provenance the
    run-log can lift verbatim, closing the v4.2 dispatch-mode declaration
    requirement.
    """
    value = arg_value or os.environ.get("INSTRUMENT_MAKER_DISPATCH_MODE")
    if value is None:
        return None
    value = value.strip().lower()
    if not value:
        return None
    if value not in VALID_DISPATCH_MODES:
        print(
            f"warning: --dispatch-mode {value!r} not in "
            f"{sorted(VALID_DISPATCH_MODES)}; stamping anyway",
            file=__import__('sys').stderr,
        )
    return value


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("packet_dir", type=Path, help="Build packet folder")
    parser.add_argument("--title", help="Presentation/packet title")
    parser.add_argument("--output-dir", type=Path, help="Output folder, defaults to packet folder")
    parser.add_argument("--images-dir", type=Path, action="append",
                        help="Extra image folder to scan; may be supplied more than once")
    parser.add_argument("--drawings-dir", type=Path, action="append",
                        help="Extra drawings folder to scan; may be supplied more than once")
    parser.add_argument(
        "--dispatch-mode",
        choices=sorted(VALID_DISPATCH_MODES),
        help="How this run was driven; stamped into capstone-manifest.json so "
             "the benchmark run log can lift it. Defaults to "
             "$INSTRUMENT_MAKER_DISPATCH_MODE if set.",
    )
    args = parser.parse_args()

    packet_dir = args.packet_dir.resolve()
    if not packet_dir.exists():
        raise SystemExit(f"Packet folder not found: {packet_dir}")

    output_dir = (args.output_dir or packet_dir).resolve()
    output_dir.mkdir(parents=True, exist_ok=True)
    title = title_from_packet(packet_dir, args.title)
    files = packet_files(packet_dir)
    slides = build_deck_slides(
        packet_dir,
        title,
        files,
        image_dirs=args.images_dir,
        drawing_dirs=args.drawings_dir,
    )

    deck_md = output_dir / "capstone-deck.md"
    print_md = output_dir / "print-packet.md"
    print_html = output_dir / "print-packet.html"
    pptx_path = output_dir / "capstone-deck.pptx"
    pdf_path = output_dir / "print-packet.pdf"
    manifest_path = output_dir / "capstone-manifest.json"

    write_deck_markdown(deck_md, slides)
    write_print_packet(print_md, packet_dir, title, files)
    print_html.write_text(markdown_to_html(print_md.read_text(encoding="utf-8"), title), encoding="utf-8")
    pptx_created = write_pptx_if_available(pptx_path, slides, title, packet_dir=packet_dir)
    pdf_created = write_pdf_if_available(pdf_path, packet_dir, title, files)

    manifest = {
        "packet_dir": str(packet_dir),
        "title": title,
        "generated": dt.date.today().isoformat(),
        "files_included": [path.name for path in files],
        "asset_discovery": {
            "images": [str(p) for p in _asset_dirs(packet_dir, "images", args.images_dir)],
            "drawings": [str(p) for p in _asset_dirs(packet_dir, "drawings", args.drawings_dir)],
        },
        "outputs": {
            "capstone_deck_markdown": str(deck_md),
            "capstone_deck_pptx": str(pptx_path) if pptx_created else None,
            "print_packet_markdown": str(print_md),
            "print_packet_html": str(print_html),
            "print_packet_pdf": str(pdf_path) if pdf_created else None,
        },
        "notes": [
            "PPTX is created when python-pptx is installed.",
            "PDF is created when reportlab (preferred) or weasyprint is installed.",
            "Use the HTML file for browser print-to-PDF when no PDF tooling is available.",
        ],
    }
    manifest_path.write_text(json.dumps(manifest, indent=2), encoding="utf-8")
    print(output_dir)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
